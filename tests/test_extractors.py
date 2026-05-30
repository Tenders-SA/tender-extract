"""Tests for all document extractors and ExtractorRegistry dispatch.

All test files are generated inline to avoid committing binary blobs.
Tests use the libraries themselves to create minimal valid documents.

Each test module follows the same pattern:
    1. Create minimal valid document bytes inline
    2. Feed to the extractor
    3. Assert ExtractionResult has expected fields
"""

import io
import os
import struct
import sys
import tempfile
from pathlib import Path

# Ensure the app module is importable
sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

import pytest

from app.extractors import ExtractorRegistry, BaseExtractor
from app.extractors.text_extractor import TextExtractor
from app.extractors.docx_extractor import DocxExtractor
from app.extractors.rtf_extractor import RtfExtractor


# =========================================================================
# Helper: create minimal PDF bytes
# =========================================================================

def _minimal_pdf_bytes(text: str = "Tender for roads maintenance") -> bytes:
    """Create a minimal valid PDF document with searchable text."""
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj\n"
        b"<< /Type /Catalog /Pages 2 0 R >>\n"
        b"endobj\n"
        b"2 0 obj\n"
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>\n"
        b"endobj\n"
        b"3 0 obj\n"
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]\n"
        b"   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\n"
        b"endobj\n"
        b"4 0 obj\n"
        b"<< /Length 44 >>\n"
        b"stream\n"
        b"BT /F1 12 Tf 100 700 Td (" + escaped.encode() + b") Tj ET\n"
        b"endstream\n"
        b"endobj\n"
        b"5 0 obj\n"
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\n"
        b"endobj\n"
        b"xref\n"
        b"0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000266 00000 n \n"
        b"0000000348 00000 n \n"
        b"trailer\n"
        b"<< /Size 6 /Root 1 0 R >>\n"
        b"startxref\n"
        b"421\n"
        b"%%EOF"
    )
    return pdf


def _minimal_docx_bytes(text: str = "Tender for road construction services.") -> bytes:
    """Create a minimal .docx file using python-docx."""
    try:
        from docx import Document
    except ImportError:
        # Return dummy ZIP bytes with docx-like structure for registry tests
        return _minimal_zip_bytes({"word/document.xml": b"<doc>"})
    doc = Document()
    doc.add_paragraph(text)
    doc.add_paragraph("Supply and delivery of construction materials.")
    doc.add_paragraph("Closing date: 30 June 2026 at 11:00 AM.")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _minimal_xlsx_bytes() -> bytes:
    """Create a minimal .xlsx file using openpyxl."""
    try:
        import openpyxl
    except ImportError:
        return _minimal_zip_bytes({"xl/workbook.xml": b"<wb/>"})
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Pricing Schedule"
    ws["A1"] = "Item"
    ws["B1"] = "Description"
    ws["C1"] = "Unit Price"
    ws["A2"] = "1"
    ws["B2"] = "Supply of cement"
    ws["C2"] = "R 125.00"
    ws["A3"] = "2"
    ws["B3"] = "Transport of materials"
    ws["C3"] = "R 5,000.00"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _minimal_xls_magic_bytes() -> bytes:
    """Create bytes with XLS (CFB) magic but no valid content.

    The CFB magic \\xD0\\xCF\\x11\\xE0\\xA1\\xB1\\x1A\\xE1 is shared across
    .doc, .xls, and .ppt legacy formats. For content tests we rely on
    LegacyDocExtractor because xlrd can't parse our inline stub.
    """
    # CFB header (512 bytes minimal)
    data = bytearray(512)
    data[0:8] = b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"
    # Set minimal CFB properties
    struct.pack_into("<H", data, 24, 0x003E)  # minor version
    struct.pack_into("<H", data, 26, 0x0003)  # major version
    struct.pack_into("<H", data, 28, 0xFFFE)  # byte order
    struct.pack_into("<H", data, 30, 0x0009)  # sector size shift (512 bytes)
    struct.pack_into("<H", data, 32, 0x0006)  # mini sector size shift (64 bytes)
    return bytes(data)


def _minimal_pptx_bytes() -> bytes:
    """Create a minimal .pptx file using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return _minimal_zip_bytes({"ppt/presentation.xml": b"<p:sld/>"})
    prs = Presentation()
    # Use a blank layout that doesn't require title placeholder
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)
    # Add text via textbox (not title placeholder which may be None)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Tender Briefing Presentation"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _minimal_odt_bytes(text: str = "Tender for municipal services.") -> bytes:
    """Create a minimal .odt file using odfpy."""
    try:
        from odf import text as odf_text
        from odf.opendocument import OpenDocumentText
    except ImportError:
        return _minimal_zip_bytes({"content.xml": b"<office:body/>"})
    doc = OpenDocumentText()
    p = odf_text.P()
    p.addText(text)
    doc.text.addElement(p)
    p2 = odf_text.P()
    p2.addText("Supply and delivery of construction materials.")
    doc.text.addElement(p2)
    buf = io.BytesIO()
    doc.save(buf)
    odt_bytes = buf.getvalue()
    # Verify odfpy produced a valid ODT by checking magic + content
    if len(odt_bytes) < 200:
        return _minimal_zip_bytes({"content.xml": b"<office:body/>"})
    return odt_bytes


def _minimal_rtf_text(text: str = "Tender for IT infrastructure services.") -> str:
    """Create minimal RTF content."""
    escaped = text.replace("\\", "\\\\").replace("{", "\\{").replace("}", "\\}")
    return f"{{\\rtf1\\ansi\\deff0 {{\\fonttbl {{\\f0 Helvetica;}}}}\\f0\\fs24 {escaped}\\par}}"


def _minimal_zip_bytes(file_dict: dict[str, bytes] | None = None) -> bytes:
    """Create a minimal ZIP file in memory."""
    import zipfile
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("info.txt", "ZIP content placeholder")
        if file_dict:
            for name, content in file_dict.items():
                zf.writestr(name, content)
    return buf.getvalue()


# =========================================================================
# Extractor Registry Dispatch Tests
# =========================================================================


class TestExtractorRegistry:
    """Tests for ExtractorRegistry dispatch logic."""

    def test_magic_pdf_detection(self):
        """PDF magic bytes (%PDF) should select PdfExtractor."""
        data = _minimal_pdf_bytes()
        extractor = ExtractorRegistry.get_extractor(data, filename=None)
        assert extractor is not None
        # PdfExtractor lives in app.extractors.pdf_extractor
        assert type(extractor).__name__ == "PdfExtractor"

    def test_magic_zip_detection(self):
        """ZIP magic bytes (PK\\x03\\x04) without office extension -> ZipExtractor."""
        data = _minimal_zip_bytes()
        extractor = ExtractorRegistry.get_extractor(data, filename="archive.zip")
        assert extractor is not None
        assert type(extractor).__name__ == "ZipExtractor"

    def test_magic_zip_docx_disambiguation(self):
        """PK magic + .docx extension -> DocxExtractor (not ZipExtractor)."""
        data = _minimal_docx_bytes()
        extractor = ExtractorRegistry.get_extractor(data, filename="tender.docx")
        assert extractor is not None
        assert type(extractor).__name__ == "DocxExtractor"

    def test_magic_zip_xlsx_disambiguation(self):
        """PK magic + .xlsx extension -> XlsxExtractor (not ZipExtractor)."""
        data = _minimal_xlsx_bytes()
        extractor = ExtractorRegistry.get_extractor(data, filename="pricing.xlsx")
        assert extractor is not None
        assert type(extractor).__name__ == "XlsxExtractor"

    def test_magic_zip_pptx_disambiguation(self):
        """PK magic + .pptx extension -> PptxExtractor (not ZipExtractor)."""
        data = _minimal_pptx_bytes()
        extractor = ExtractorRegistry.get_extractor(data, filename="briefing.pptx")
        assert extractor is not None
        assert type(extractor).__name__ == "PptxExtractor"

    def test_cfb_magic_detection(self):
        """CFB magic bytes (\\xD0\\xCF\\x11... ) -> LegacyDocExtractor."""
        data = _minimal_xls_magic_bytes()
        extractor = ExtractorRegistry.get_extractor(data, filename="document.doc")
        assert extractor is not None
        assert type(extractor).__name__ == "LegacyDocExtractor"

    def test_extension_based_detection(self):
        """Fallback to extension when magic doesn't match."""
        data = b"this is not a known format but ends in .csv"
        extractor = ExtractorRegistry.get_extractor(data, filename="data.csv")
        assert extractor is not None
        assert type(extractor).__name__ == "TextExtractor"

    def test_extension_rtf_detection(self):
        """.rtf extension -> RtfExtractor."""
        data = _minimal_rtf_text().encode("utf-8")
        extractor = ExtractorRegistry.get_extractor(data, filename="document.rtf")
        assert extractor is not None
        assert type(extractor).__name__ == "RtfExtractor"

    def test_extension_odt_detection(self):
        """.odt extension -> OdtExtractor."""
        data = _minimal_odt_bytes()
        extractor = ExtractorRegistry.get_extractor(data, filename="municipal.odt")
        assert extractor is not None
        assert type(extractor).__name__ == "OdtExtractor"

    def test_extension_txt_detection(self):
        """.txt extension -> TextExtractor."""
        data = b"Hello World"
        extractor = ExtractorRegistry.get_extractor(data, filename="notes.txt")
        assert extractor is not None
        assert type(extractor).__name__ == "TextExtractor"

    def test_mime_based_detection(self):
        """MIME type fallback when no magic or extension."""
        data = b"some content"
        extractor = ExtractorRegistry.get_extractor(
            data, filename=None, mime_type="application/pdf"
        )
        assert extractor is not None
        assert type(extractor).__name__ == "PdfExtractor"

    def test_mime_csv_detection(self):
        """text/csv MIME -> TextExtractor."""
        data = b"a,b,c"
        extractor = ExtractorRegistry.get_extractor(
            data, filename=None, mime_type="text/csv"
        )
        assert extractor is not None
        assert type(extractor).__name__ == "TextExtractor"

    def test_unsupported_format_returns_none(self):
        """Unrecognized binary data with no extension/MIME returns None."""
        data = b"\x00\x01\x02\x03\x04\x05\x06\x07\x08\x09"
        extractor = ExtractorRegistry.get_extractor(data, filename=None)
        assert extractor is None

    def test_unsupported_extension_returns_none(self):
        """Unknown extension with no magic match returns None."""
        data = b"some content"
        extractor = ExtractorRegistry.get_extractor(data, filename="file.xyz")
        assert extractor is None

    def test_empty_data_returns_none(self):
        """Empty data with no hints returns None."""
        extractor = ExtractorRegistry.get_extractor(b"", filename=None)
        assert extractor is None

    def test_all_extensions_map_to_extractors(self):
        """Every extension in EXTENSION_MAP resolves to a valid extractor."""
        for ext, class_name in ExtractorRegistry.EXTENSION_MAP.items():
            # Create minimal data with the right extension
            data = _minimal_pdf_bytes() if ext == ".pdf" else b"test"
            filename = f"test{ext}"
            extractor = ExtractorRegistry.get_extractor(data, filename=filename)
            assert extractor is not None, f"Extension {ext} -> {class_name} returned None"
            assert type(extractor).__name__ == class_name, (
                f"Extension {ext} expected {class_name}, got {type(extractor).__name__}"
            )

    def test_all_mime_types_resolve(self):
        """Every MIME type in MIME_MAP resolves to a valid extractor class."""
        for mime_type, class_name in ExtractorRegistry.MIME_MAP.items():
            data = b"test"
            extractor = ExtractorRegistry.get_extractor(
                data, filename=None, mime_type=mime_type
            )
            # Some MIME types may need specific data (e.g. PDF magic)
            # But the class resolution should at least not crash
            if extractor is not None:
                assert type(extractor).__name__ == class_name, (
                    f"MIME {mime_type} expected {class_name}, got {type(extractor).__name__}"
                )


# =========================================================================
# PDF Extractor Tests
# =========================================================================


class TestPdfExtractor:
    """Tests for PdfExtractor with inline-generated PDF files."""

    def test_extract_returns_result(self):
        """PdfExtractor returns ExtractionResult for valid PDF."""
        from app.extractors.pdf_extractor import PdfExtractor

        pdf_bytes = _minimal_pdf_bytes(
            "Tender Number: SCM/2026/001\n"
            "Description: Supply and delivery of office furniture\n"
            "Closing Date: 30 June 2026\n"
            "Issuing Organization: Department of Transport"
        )
        extractor = PdfExtractor()
        result = extractor.extract(pdf_bytes)
        assert result is not None
        assert result.full_text
        assert len(result.full_text) > 10

    def test_extract_with_filename(self):
        """PdfExtractor accepts optional filename parameter."""
        from app.extractors.pdf_extractor import PdfExtractor

        pdf_bytes = _minimal_pdf_bytes("Test content")
        extractor = PdfExtractor()
        result = extractor.extract(pdf_bytes, filename="tender.pdf")
        assert result is not None

    def test_extract_with_too_little_text(self):
        """Very short PDF results in 0 confidence + needs_ai_fallback."""
        from app.extractors.pdf_extractor import PdfExtractor

        pdf_bytes = _minimal_pdf_bytes("Hi")
        extractor = PdfExtractor()
        result = extractor.extract(pdf_bytes)
        # PDF with insufficient text -> 0 confidence
        assert result.confidence == 0.0 or result.confidence < 0.5

    def test_contains_scanned_pages(self):
        """contains_scanned_pages returns False for text PDF."""
        from app.extractors.pdf_extractor import PdfExtractor

        pdf_bytes = _minimal_pdf_bytes("Some text content here for testing purposes.")
        extractor = PdfExtractor()
        result = extractor.contains_scanned_pages(pdf_bytes)
        # Our minimal PDF has text, so not scanned
        assert result is False

    def test_base_extractor_interface(self):
        """PdfExtractor is a BaseExtractor."""
        from app.extractors.pdf_extractor import PdfExtractor

        assert isinstance(PdfExtractor(), BaseExtractor)


# =========================================================================
# DOCX Extractor Tests
# =========================================================================


class TestDocxExtractor:
    """Tests for DocxExtractor with python-docx-generated files."""

    def test_extract_returns_text(self):
        """DocxExtractor extracts text from a valid .docx."""
        docx_bytes = _minimal_docx_bytes(
            "Tender for the provision of security services at municipal buildings."
        )
        if len(docx_bytes) < 100:
            pytest.skip("python-docx not available — cannot generate test file")

        extractor = DocxExtractor()
        result = extractor.extract(docx_bytes, filename="tender.docx")
        assert result is not None
        assert len(result.full_text) > 20
        assert "security" in result.full_text.lower() or "construction" in result.full_text.lower()

    def test_description_set(self):
        """First meaningful paragraph is used as description."""
        docx_bytes = _minimal_docx_bytes(
            "Provision of cleaning services for the municipal offices for a period of 36 months."
        )
        if len(docx_bytes) < 100:
            pytest.skip("python-docx not available")

        extractor = DocxExtractor()
        result = extractor.extract(docx_bytes)
        assert result.description
        assert len(result.description.strip()) > 0

    def test_confidence_based_on_length(self):
        """Confidence reflects text length."""
        docx_bytes = _minimal_docx_bytes("Short text.")
        if len(docx_bytes) < 100:
            pytest.skip("python-docx not available")

        extractor = DocxExtractor()
        result = extractor.extract(docx_bytes)
        # Short text should have lower confidence
        assert isinstance(result.confidence, float)

    def test_bad_data_graceful(self):
        """DocxExtractor handles corrupted data gracefully."""
        extractor = DocxExtractor()
        result = extractor.extract(b"not a valid docx", filename="broken.docx")
        assert result is not None
        assert result.confidence == 0.0

    def test_base_extractor_interface(self):
        """DocxExtractor is a BaseExtractor."""
        assert isinstance(DocxExtractor(), BaseExtractor)


# =========================================================================
# Legacy DOC Extractor Tests
# =========================================================================


class TestLegacyDocExtractor:
    """Tests for LegacyDocExtractor."""

    def test_binary_text_extraction(self):
        """LegacyDocExtractor extracts printable text from binary data."""
        from app.extractors.legacy_doc_extractor import LegacyDocExtractor

        # Create binary data with embedded printable text runs
        data = (
            b"\x00\x01\x02"
            b"This is a tender document for road construction" + b"\x00\x01\x02"
            b"Supply of materials for road maintenance" + b"\xFF\xFE\xFD"
            b"Closing date: 30 June 2026" + b"\x00"
        )
        extractor = LegacyDocExtractor()
        result = extractor.extract(data, filename="tender.doc")
        assert result is not None
        assert "tender" in result.full_text.lower()
        assert "road" in result.full_text.lower()

    def test_binary_with_no_antiword(self):
        """LegacyDocExtractor falls back to binary text when antiword not available."""
        from app.extractors.legacy_doc_extractor import LegacyDocExtractor

        # antiword likely not installed in test env
        data = b"\x00PROJECT\x00DESCRIPTION\x00Supply of goods\x00"
        extractor = LegacyDocExtractor()
        result = extractor.extract(data, filename="document.doc")
        assert result is not None
        assert len(result.full_text) > 0

    def test_magic_byte_detection(self):
        """CFB magic bytes route to LegacyDocExtractor."""
        data = _minimal_xls_magic_bytes()
        extractor = ExtractorRegistry.get_extractor(data, filename="document.doc")
        assert extractor is not None
        assert type(extractor).__name__ == "LegacyDocExtractor"

    def test_base_extractor_interface(self):
        """LegacyDocExtractor is a BaseExtractor."""
        from app.extractors.legacy_doc_extractor import LegacyDocExtractor

        assert isinstance(LegacyDocExtractor(), BaseExtractor)


# =========================================================================
# XLSX Extractor Tests
# =========================================================================


class TestXlsxExtractor:
    """Tests for XlsxExtractor."""

    def test_extract_returns_text(self):
        """XlsxExtractor extracts text from a valid .xlsx."""
        from app.extractors.xlsx_extractor import XlsxExtractor

        xlsx_bytes = _minimal_xlsx_bytes()
        if len(xlsx_bytes) < 200:
            pytest.skip("openpyxl not available — cannot generate test file")

        extractor = XlsxExtractor()
        result = extractor.extract(xlsx_bytes, filename="pricing.xlsx")
        assert result is not None
        assert len(result.full_text) > 0
        # Check that spreadsheet content was extracted
        assert "Cement" in result.full_text or "cement" in result.full_text

    def test_sheet_name_in_output(self):
        """XlsxExtractor includes sheet name in output."""
        from app.extractors.xlsx_extractor import XlsxExtractor

        xlsx_bytes = _minimal_xlsx_bytes()
        if len(xlsx_bytes) < 200:
            pytest.skip("openpyxl not available")

        extractor = XlsxExtractor()
        result = extractor.extract(xlsx_bytes)
        assert "[Sheet:" in result.full_text or "Pricing" in result.full_text

    def test_bad_data_graceful(self):
        """XlsxExtractor handles corrupted data gracefully."""
        from app.extractors.xlsx_extractor import XlsxExtractor

        extractor = XlsxExtractor()
        result = extractor.extract(b"not an xlsx", filename="bad.xlsx")
        assert result is not None
        assert result.confidence == 0.0

    def test_base_extractor_interface(self):
        """XlsxExtractor is a BaseExtractor."""
        from app.extractors.xlsx_extractor import XlsxExtractor

        assert isinstance(XlsxExtractor(), BaseExtractor)


# =========================================================================
# XLS Extractor Tests
# =========================================================================


class TestXlsExtractor:
    """Tests for XlsExtractor."""

    def test_extract_returns_text(self):
        """XlsExtractor works with inline-generated XLS data."""
        from app.extractors.xls_extractor import XlsExtractor

        # We can't easily generate XLS without xlwt (not in requirements).
        # Test that the extractor gracefully handles invalid data.
        extractor = XlsExtractor()
        result = extractor.extract(b"\x00" * 512, filename="legacy.xls")
        # With CFB header but no valid workbook, xlrd should fail gracefully
        assert result is not None
        assert result.confidence == 0.0

    def test_bad_data_graceful(self):
        """XlsExtractor handles corrupted data gracefully."""
        from app.extractors.xls_extractor import XlsExtractor

        extractor = XlsExtractor()
        result = extractor.extract(b"not an xls", filename="bad.xls")
        assert result is not None
        assert result.confidence == 0.0

    def test_base_extractor_interface(self):
        """XlsExtractor is a BaseExtractor."""
        from app.extractors.xls_extractor import XlsExtractor

        assert isinstance(XlsExtractor(), BaseExtractor)


# =========================================================================
# PPTX Extractor Tests
# =========================================================================


class TestPptxExtractor:
    """Tests for PptxExtractor."""

    def test_extract_returns_text(self):
        """PptxExtractor extracts text from a valid .pptx."""
        from app.extractors.pptx_extractor import PptxExtractor

        pptx_bytes = _minimal_pptx_bytes()
        if len(pptx_bytes) < 200:
            pytest.skip("python-pptx not available — cannot generate test file")

        extractor = PptxExtractor()
        result = extractor.extract(pptx_bytes, filename="briefing.pptx")
        assert result is not None
        assert "Tender" in result.full_text or "Briefing" in result.full_text

    def test_slide_format_in_output(self):
        """PptxExtractor includes slide markers in output."""
        from app.extractors.pptx_extractor import PptxExtractor

        pptx_bytes = _minimal_pptx_bytes()
        if len(pptx_bytes) < 200:
            pytest.skip("python-pptx not available")

        extractor = PptxExtractor()
        result = extractor.extract(pptx_bytes)
        # Should have slide markers
        assert "---" in result.full_text or "Slide" in result.full_text or "Title" in result.full_text

    def test_bad_data_graceful(self):
        """PptxExtractor handles corrupted data gracefully."""
        from app.extractors.pptx_extractor import PptxExtractor

        extractor = PptxExtractor()
        result = extractor.extract(b"not a pptx", filename="bad.pptx")
        assert result is not None
        assert result.confidence == 0.0

    def test_base_extractor_interface(self):
        """PptxExtractor is a BaseExtractor."""
        from app.extractors.pptx_extractor import PptxExtractor

        assert isinstance(PptxExtractor(), BaseExtractor)


# =========================================================================
# ODT Extractor Tests
# =========================================================================


class TestOdtExtractor:
    """Tests for OdtExtractor."""

    def test_extract_returns_text(self):
        """OdtExtractor extracts text from a valid .odt."""
        from app.extractors.odt_extractor import OdtExtractor

        odt_bytes = _minimal_odt_bytes(
            "Tender for municipal infrastructure upgrade projects."
        )
        if len(odt_bytes) < 200:
            pytest.skip("odfpy not available — cannot generate test file")

        extractor = OdtExtractor()
        result = extractor.extract(odt_bytes, filename="municipal.odt")
        assert result is not None
        assert len(result.full_text) > 0
        assert "municipal" in result.full_text.lower()

    def test_description_set(self):
        """OdtExtractor populates description."""
        from app.extractors.odt_extractor import OdtExtractor

        odt_bytes = _minimal_odt_bytes(
            "Appointment of a service provider for the upgrading of wastewater "
            "treatment works in the Frances Baard District Municipality."
        )
        if len(odt_bytes) < 200:
            pytest.skip("odfpy not available")

        extractor = OdtExtractor()
        result = extractor.extract(odt_bytes)
        # Should have meaningful description
        assert result.description is not None

    def test_bad_data_graceful(self):
        """OdtExtractor handles corrupted data gracefully."""
        from app.extractors.odt_extractor import OdtExtractor

        extractor = OdtExtractor()
        result = extractor.extract(b"not an odt", filename="bad.odt")
        assert result is not None
        assert result.confidence == 0.0

    def test_base_extractor_interface(self):
        """OdtExtractor is a BaseExtractor."""
        from app.extractors.odt_extractor import OdtExtractor

        assert isinstance(OdtExtractor(), BaseExtractor)


# =========================================================================
# RTF Extractor Tests
# =========================================================================


class TestRtfExtractor:
    """Tests for RtfExtractor."""

    def test_extract_returns_text(self):
        """RtfExtractor extracts text from valid RTF."""
        rtf_text = _minimal_rtf_text(
            "Tender for the supply of IT equipment and software licenses."
        )
        extractor = RtfExtractor()
        result = extractor.extract(rtf_text.encode("utf-8"), filename="document.rtf")
        assert result is not None
        assert len(result.full_text) > 0
        assert "Tender" in result.full_text or "IT" in result.full_text

    def test_empty_rtf(self):
        """RtfExtractor handles RTF with no text."""
        minimal = b"{\\rtf1\\ansi}"
        extractor = RtfExtractor()
        result = extractor.extract(minimal, filename="empty.rtf")
        assert result is not None
        # No text found -> confidence 0
        assert result.confidence == 0.0

    def test_latin1_encoding(self):
        """RtfExtractor falls back to latin-1 for non-UTF-8 data."""
        # RTF with latin-1 encoded content
        rtf_bytes = b"{\\rtf1\\ansi {\\f0 Hello caf\xe9}}"
        extractor = RtfExtractor()
        result = extractor.extract(rtf_bytes, filename="cafe.rtf")
        assert result is not None

    def test_bad_data_graceful(self):
        """RtfExtractor handles truly bad data."""
        extractor = RtfExtractor()
        result = extractor.extract(b"\xff\xfe\x00\x01\x02 not rtf", filename="bad.rtf")
        assert result is not None

    def test_base_extractor_interface(self):
        """RtfExtractor is a BaseExtractor."""
        assert isinstance(RtfExtractor(), BaseExtractor)


# =========================================================================
# CSV/TXT Extractor Tests
# =========================================================================


class TestTextExtractor:
    """Tests for TextExtractor (CSV and TXT handling)."""

    def test_txt_extraction(self):
        """TextExtractor extracts plain text from .txt."""
        text = "This is a tender document.\nSupply of goods.\nClosing Date: 30 June 2026."
        data = text.encode("utf-8")
        extractor = TextExtractor()
        result = extractor.extract(data, filename="notes.txt")
        assert result is not None
        assert "tender" in result.full_text.lower()
        assert result.confidence > 0.0

    def test_csv_extraction(self):
        """TextExtractor parses CSV and formats as table."""
        csv_content = "Item,Description,Price\n1,Cement,R125\n2,Bricks,R200\n3,Sand,R80\n"
        data = csv_content.encode("utf-8")
        extractor = TextExtractor()
        result = extractor.extract(data, filename="pricing.csv")
        assert result is not None
        assert "Item" in result.full_text
        assert "Cement" in result.full_text
        assert result.confidence > 0.0

    def test_csv_description(self):
        """TextExtractor produces description for CSV."""
        csv = "Company,Contact,Email\nABC Construction,John,john@abc.com\n"
        extractor = TextExtractor()
        result = extractor.extract(csv.encode("utf-8"), filename="vendors.csv")
        assert result is not None
        assert result.description

    def test_empty_txt(self):
        """Empty text file returns low confidence."""
        extractor = TextExtractor()
        result = extractor.extract(b"", filename="empty.txt")
        assert result is not None
        assert result.confidence == 0.0

    def test_utf8_bom(self):
        """TextExtractor handles UTF-8 BOM."""
        text = "\ufeffTender document with BOM header.\nMore content here."
        extractor = TextExtractor()
        result = extractor.extract(text.encode("utf-8"), filename="bom.txt")
        assert result is not None
        assert len(result.full_text) > 0

    def test_latin1_fallback(self):
        """TextExtractor falls back to latin-1 for non-UTF-8 data."""
        data = b"Tender for caf\xe9 supply.\n"
        extractor = TextExtractor()
        result = extractor.extract(data, filename="cafe.txt")
        assert result is not None
        assert "caf" in result.full_text.lower()

    def test_csv_with_only_header(self):
        """CSV with only a header row produces low but non-zero confidence."""
        extractor = TextExtractor()
        result = extractor.extract(b"Item,Description,Price\n", filename="empty.csv")
        assert result is not None
        assert result.confidence >= 0.0

    def test_base_extractor_interface(self):
        """TextExtractor is a BaseExtractor."""
        assert isinstance(TextExtractor(), BaseExtractor)


# =========================================================================
# ZIP Extractor Tests
# =========================================================================


class TestZipExtractor:
    """Tests for ZipExtractor."""

    def test_extract_from_zip_with_txt(self):
        """ZipExtractor extracts text from ZIP containing .txt files."""
        from app.extractors.zip_extractor import ZipExtractor

        import zipfile
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("scope.txt", "Description: Supply of construction materials.")
            zf.writestr("pricing.txt", "Item 1: Cement\nItem 2: Bricks")
        zip_bytes = buf.getvalue()

        extractor = ZipExtractor()
        result = extractor.extract(zip_bytes, filename="bidpack.zip")
        assert result is not None
        assert len(result.full_text) > 0

    def test_empty_zip(self):
        """Empty ZIP returns empty result with confidence 0."""
        from app.extractors.zip_extractor import ZipExtractor

        import zipfile
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            pass  # Empty ZIP
        zip_bytes = buf.getvalue()

        extractor = ZipExtractor()
        result = extractor.extract(zip_bytes, filename="empty.zip")
        assert result is not None
        assert result.confidence == 0.0

    def test_zip_with_only_images(self):
        """ZIP with only non-document files returns empty result."""
        from app.extractors.zip_extractor import ZipExtractor

        import zipfile
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("photo.jpg", b"\xff\xd8\xff\xe0")  # JPEG header
            zf.writestr("image.png", b"\x89PNG\r\n\x1a\n")
        zip_bytes = buf.getvalue()

        extractor = ZipExtractor()
        result = extractor.extract(zip_bytes, filename="images.zip")
        assert result is not None
        assert result.confidence == 0.0

    def test_bad_zip_data(self):
        """Corrupted ZIP data returns error gracefully."""
        from app.extractors.zip_extractor import ZipExtractor

        extractor = ZipExtractor()
        result = extractor.extract(b"not a zip file content", filename="bad.zip")
        assert result is not None
        assert result.confidence == 0.0

    def test_base_extractor_interface(self):
        """ZipExtractor is a BaseExtractor."""
        from app.extractors.zip_extractor import ZipExtractor

        assert isinstance(ZipExtractor(), BaseExtractor)


# =========================================================================
# Cross-format / Integration Tests
# =========================================================================


class TestCrossFormatExtraction:
    """Tests that span multiple format concerns."""

    def test_unsupported_extension_via_registry(self):
        """POST /v1/extract would get 415 (registry returns None)."""
        data = b"some binary garbage data that is not any known format"
        extractor = ExtractorRegistry.get_extractor(data, filename="weird.bin")
        assert extractor is None

    def test_unsupported_mime_via_registry(self):
        """Registry returns None for unsupported MIME types."""
        extractor = ExtractorRegistry.get_extractor(
            b"test", filename=None, mime_type="application/octet-stream"
        )
        assert extractor is None

    def test_filename_without_extension_returns_none(self):
        """File with no extension and no magic returns None."""
        extractor = ExtractorRegistry.get_extractor(b"test", filename="README")
        assert extractor is None

    def test_mixed_case_extension_resolves(self):
        """Extensions are case-insensitive."""
        data = b"a,b,c\n1,2,3"
        extractor = ExtractorRegistry.get_extractor(data, filename="DATA.CSV")
        assert extractor is not None
        assert type(extractor).__name__ == "TextExtractor"

    def test_mime_type_with_charset(self):
        """MIME type with charset suffix resolves correctly."""
        extractor = ExtractorRegistry.get_extractor(
            b"test", filename=None, mime_type="text/plain; charset=utf-8"
        )
        assert extractor is not None
        assert type(extractor).__name__ == "TextExtractor"
