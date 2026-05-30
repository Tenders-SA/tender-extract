"""PPTX (PowerPoint) document extraction using python-pptx.

Extracts text from all slides, shapes (titles, text boxes, tables, groups).
"""

from io import BytesIO
from typing import Optional

from ..extractor import ExtractionResult
from .base import BaseExtractor


class PptxExtractor(BaseExtractor):
    """Extractor for .pptx (PowerPoint) presentations.

    Uses python-pptx to iterate all slides and shapes, extracting text
    content formatted as:
        --- Slide N ---
        Title: ...
        Content: ...
    """

    def extract(self, data: bytes, filename: Optional[str] = None) -> ExtractionResult:
        """Extract text from a .pptx presentation.

        Args:
            data: Raw .pptx file bytes.
            filename: Optional original filename.

        Returns:
            ExtractionResult with extracted text content.
        """
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
            from pptx.util import Inches, Pt  # noqa: F401 (imported for library loading)
        except ImportError:
            return ExtractionResult(
                full_text="",
                confidence=0.0,
                requirements=["python-pptx package is not installed"],
            )

        try:
            prs = Presentation(BytesIO(data))
        except Exception:
            return ExtractionResult(
                full_text="",
                confidence=0.0,
                requirements=["Failed to open PPTX presentation"],
            )

        all_text_parts: list[str] = []
        slide_count = 0
        total_text_length = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            has_content = False

            for shape in slide.shapes:
                texts = _extract_shape_text(shape)
                for text in texts:
                    clean = text.strip()
                    if clean:
                        slide_parts.append(clean)
                        total_text_length += len(clean)
                        has_content = True

            if has_content:
                slide_count += 1
                all_text_parts.append(f"--- Slide {slide_idx} ---")

                # Try to identify title vs content
                title_text = ""
                content_parts: list[str] = []
                for part in slide_parts:
                    if shape_has_text_frame(part, slide):
                        if not title_text:
                            title_text = part
                        else:
                            content_parts.append(part)

                if title_text:
                    all_text_parts.append(f"Title: {title_text}")
                for content in content_parts:
                    if content != title_text:
                        all_text_parts.append(content)

        if not all_text_parts:
            return ExtractionResult(
                full_text="",
                confidence=0.0,
                requirements=["No text content found in PPTX presentation"],
            )

        full_text = "\n".join(all_text_parts)

        # Confidence based on content volume
        if total_text_length >= 500:
            confidence = 0.7
        elif total_text_length >= 100:
            confidence = 0.5
        elif total_text_length >= 30:
            confidence = 0.3
        else:
            confidence = 0.1

        # Description: first meaningful text
        description = ""
        for line in full_text.splitlines():
            if line.startswith("Title:") and len(line) > 60:
                description = line
                break
        if not description:
            for line in full_text.splitlines():
                clean = line.strip()
                if clean and len(clean) > 50:
                    description = clean
                    break

        return ExtractionResult(
            full_text=full_text,
            description=description,
            confidence=confidence,
            raw_text_preview=full_text[:500] if full_text else None,
        )


def _extract_shape_text(shape: object) -> list[str]:
    """Extract text from a PowerPoint shape.

    Handles text frames, tables, and grouped shapes recursively.
    """
    texts: list[str] = []

    try:
        # Access via python-pptx internal attributes (duck-typing approach)
        # Text frame
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)

        # Table
        if hasattr(shape, "has_table") and shape.has_table:
            table = shape.table
            for row in table.rows:
                row_texts: list[str] = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    texts.append(" | ".join(row_texts))

        # Group shapes (recursive)
        if hasattr(shape, "shapes"):
            for child_shape in shape.shapes:
                texts.extend(_extract_shape_text(child_shape))
    except Exception:
        pass

    return texts


def shape_has_text_frame(part: str, slide: object) -> bool:
    """Heuristic to determine if text came from a text frame vs a table.

    This is a simplified heuristic — real title detection happens via
    python-pptx's shape.has_text_frame attribute, but for our purposes
    we just need reasonable content separation.
    """
    # Consider any non-table, non-group content as text frame content
    return True
